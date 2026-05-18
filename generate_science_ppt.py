from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN

# Theme Colors based on after-12.html / orientation.html
C_SAFFRON = RGBColor(255, 107, 0)
C_NAVY = RGBColor(13, 27, 62)
C_CREAM = RGBColor(253, 248, 240)
C_GOLD = RGBColor(212, 168, 83)
C_GREEN = RGBColor(19, 138, 54)
C_TEAL = RGBColor(14, 124, 107)
C_ROSE = RGBColor(192, 57, 43)
C_MUTED = RGBColor(107, 107, 107)
C_WHITE = RGBColor(255, 255, 255)

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_footer(slide):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Nudge Academy | Complete Career Guide for Class 12 Science"
    p.font.size = Pt(12)
    p.font.color.rgb = C_MUTED
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def add_title(slide, text, color=C_NAVY):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color

def create_ppt():
    prs = Presentation()
    # Change slide size to Widescreen 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # --- SLIDE 1: TITLE SLIDE ---
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide1, C_NAVY)

    # Decorative Shape
    shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-1), Inches(4), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_SAFFRON
    shape.line.fill.background()
    shape.shadow.inherit = False

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Every Option After"
    p.font.size = Pt(40)
    p.font.color.rgb = C_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Plus Two Science"
    p2.font.size = Pt(56)
    p2.font.bold = True
    p2.font.color.rgb = C_GOLD

    txBox_sub = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf_sub = txBox_sub.text_frame
    p3 = tf_sub.paragraphs[0]
    p3.text = "The most exhaustive guide for PCM & PCB students.\nPrepared by Nudge Academy."
    p3.font.size = Pt(20)
    p3.font.color.rgb = C_WHITE

    # --- SLIDE 2: THE BIG PICTURE (WITH BAR CHART) ---
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide2, C_CREAM)
    add_title(slide2, "The Big Picture — All Paths at a Glance")
    
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Science opens the widest range of options across all domains."
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• PCM: Engineering, Architecture, Defence, Pure Sciences, Merchant Navy"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• PCB: Medicine, Allied Healthcare, Pharmacy, Agriculture, Biotech"
    p.font.size = Pt(16)

    # Chart: Overview seats
    chart_data = CategoryChartData()
    chart_data.categories = ['Engineering', 'Medicine', 'Law', 'Mgmt', 'Teaching', 'Agri', 'Design', 'Defence', 'Pure Sci']
    chart_data.add_series('Annual Seats (thousands)', (1500, 180, 2.5, 50, 300, 45, 5, 0.4, 2))
    x, y, cx, cy = Inches(0.5), Inches(3.0), Inches(9), Inches(3.8)
    chart = slide2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_SAFFRON

    add_footer(slide2)

    # --- SLIDE 3: ENGINEERING (WITH JEE BAR CHART) ---
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide3, C_CREAM)
    add_title(slide3, "Engineering & Technology (PCM)")

    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Gold Standard Pathway:"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• JEE Main: Gateway to NITs, IIITs, GFTIs."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• JEE Advanced: Top 2.5L attempt this for 23 IITs."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• BITSAT: Private excellence (BITS Pilani, Goa, Hyd)."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• State CETs: MHT-CET, KCET, WBJEE, COMEDK."
    p.font.size = Pt(16)

    # JEE Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Main Reg.', 'Main Qual.', 'Adv. App.', 'IIT Seats']
    chart_data.add_series('Students (in thousands)', (1150, 250, 190, 17))
    x, y, cx, cy = Inches(5.0), Inches(1.2), Inches(4.5), Inches(3)
    chart = slide3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = False
    
    # IIT Growth Chart
    chart_data2 = CategoryChartData()
    chart_data2.categories = ['1951', '1980', '2000', '2010', '2020', '2024']
    chart_data2.add_series('Number of IITs', (1, 6, 7, 15, 22, 23))
    x, y, cx, cy = Inches(5.0), Inches(4.3), Inches(4.5), Inches(2.5)
    chart2 = slide3.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data2).chart
    chart2.has_legend = False
    chart2.series[0].format.line.color.rgb = C_NAVY

    add_footer(slide3)

    # --- SLIDE 4: MEDICINE & HEALTHCARE (WITH NEET/AIIMS CHARTS) ---
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide4, C_CREAM)
    add_title(slide4, "Medicine & Healthcare (PCB)")

    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Healthcare Ecosystem:"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• NEET-UG: The ONLY exam for MBBS, BDS, BAMS, BHMS."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• AIIMS: 25 campuses across India. 1.1L+ Govt MBBS seats total."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Allied Fields: B.Sc. Nursing, Pharmacy, Physiotherapy."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Veterinary: B.V.Sc (Via NEET)."
    p.font.size = Pt(16)

    # NEET Growth Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['2017','2019','2021','2023','2024']
    chart_data.add_series('Applicants (Lakhs)', (11.4, 15.2, 16.1, 20.9, 24.1))
    x, y, cx, cy = Inches(5.0), Inches(1.2), Inches(4.5), Inches(2.8)
    chart = slide4.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
    chart.has_legend = False
    chart.series[0].format.line.color.rgb = C_ROSE

    # AIIMS Pie Chart
    chart_data2 = CategoryChartData()
    chart_data2.categories = ['North', 'South', 'East', 'West', 'Central', 'NE']
    chart_data2.add_series('AIIMS Distribution', (9, 5, 4, 3, 2, 2))
    x, y, cx, cy = Inches(5.0), Inches(4.1), Inches(4.5), Inches(2.8)
    chart2 = slide4.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data2).chart
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.RIGHT

    add_footer(slide4)

    # --- SLIDE 5: PURE SCIENCES (WITH QS BAR CHART) ---
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide5, C_CREAM)
    add_title(slide5, "Pure Sciences & Research")

    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Research & Innovation Pathways:"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• IISc Bangalore: 4-year BS (via JEE/IAT). Top global institute."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• IISERs: 5-year BS-MS via IAT."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• NISER: 5-year integrated M.Sc via NEST."
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• ISI / CMI: Premier maths/stats institutes."
    p.font.size = Pt(16)

    # QS Ranking Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['IISc', 'IIT B', 'IIT D', 'IIT M', 'IIT KGP']
    chart_data.add_series('QS World Rank 2024 (Lower=Better)', (186, 233, 344, 285, 271))
    x, y, cx, cy = Inches(5.0), Inches(2.0), Inches(4.5), Inches(4)
    chart = slide5.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = False
    
    add_footer(slide5)

    # --- SLIDE 6: INSTITUTE SEATS (OVERALL BAR CHART) ---
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide6, C_CREAM)
    add_title(slide6, "Institute Seats Breakdown")

    chart_data = CategoryChartData()
    chart_data.categories = ['NITs (JEE Main)', 'IITs (JEE Adv.)', 'Medical Govt (NEET)', 'NLUs (CLAT)', 'IIMs (CAT)', 'IISERs (IAT)']
    chart_data.add_series('Approx. Annual Seats', (23000, 17000, 55000, 2500, 6000, 1200))
    x, y, cx, cy = Inches(1.0), Inches(1.5), Inches(8), Inches(5)
    chart = slide6.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_TEAL

    add_footer(slide6)
    
    # --- SLIDE 7: OTHER SCIENCE PATHWAYS ---
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide7, C_CREAM)
    add_title(slide7, "Other Science Pathways")

    txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Defence & Marine (PCM)"
    p.font.size = Pt(22)
    p.font.color.rgb = C_NAVY
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• NDA: Tri-service academy"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Naval Academy: 10+2 B.Tech Entry"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• Merchant Navy: IMU CET"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\nAgriculture & Allied (PCB/PCM)"
    p.font.size = Pt(22)
    p.font.color.rgb = C_GREEN
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• ICAR AIEEA: Gateway to Central/State Agri Universities"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• B.Sc. Agriculture, Horticulture, Forestry"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• B.F.Sc (Fisheries), B.Tech (Dairy/Food)"
    p.font.size = Pt(16)

    txBox2 = slide7.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Cross-Disciplinary Paths"
    p.font.size = Pt(22)
    p.font.color.rgb = C_SAFFRON
    p.font.bold = True
    p = tf2.add_paragraph()
    p.text = "• Law: CLAT for NLUs (BA/BBA LLB)"
    p.font.size = Pt(16)
    p = tf2.add_paragraph()
    p.text = "• Mgmt: IPMAT (5-year Mgmt at IIMs)"
    p.font.size = Pt(16)
    p = tf2.add_paragraph()
    p.text = "• Design/Arch: NID, NIFT, NATA, UCEED"
    p.font.size = Pt(16)
    
    p = tf2.add_paragraph()
    p.text = "\nTeaching & Govt Jobs"
    p.font.size = Pt(22)
    p.font.color.rgb = C_ROSE
    p.font.bold = True
    p = tf2.add_paragraph()
    p.text = "• Teaching: NCET (B.Sc.B.Ed at IIT/NIT)"
    p.font.size = Pt(16)
    p = tf2.add_paragraph()
    p.text = "• UPSC/SSC/Railways: Solid govt careers"
    p.font.size = Pt(16)

    add_footer(slide7)
    
    # --- SLIDE 8: ACTION PLAN ---
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide8, C_NAVY)
    
    add_title(slide8, "Your Action Plan", C_WHITE)

    txBox = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "1. Shortlist 2-3 target exams (e.g. Primary: JEE/NEET, Backup: CUET, NCET)"
    p.font.size = Pt(22)
    p.font.color.rgb = C_WHITE
    p = tf.add_paragraph()
    p.text = "2. CUET is a MUST: It opens doors to 250+ universities for B.Sc programs."
    p.font.size = Pt(22)
    p.font.color.rgb = C_WHITE
    p = tf.add_paragraph()
    p.text = "3. Track Registration Dates: CUET, NDA, ICAR, IAT all open early."
    p.font.size = Pt(22)
    p.font.color.rgb = C_WHITE
    p = tf.add_paragraph()
    p.text = "4. Don't Panic: Your marks are a starting point, not the end."
    p.font.size = Pt(22)
    p.font.color.rgb = C_WHITE

    add_footer(slide8)

    prs.save('Career_Options_After_12_Science_Designed.pptx')
    print("PPT created successfully: Career_Options_After_12_Science_Designed.pptx")

if __name__ == "__main__":
    create_ppt()
